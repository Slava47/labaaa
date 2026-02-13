#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Лабораторная работа №8
Microsoft Office PowerPoint. Создание презентации
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_lab8_presentation():
    """Создание презентации PowerPoint"""
    
    # Создание презентации
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Слайд 1: Титульный
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "О МОЕМ ДРУГЕ"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle.text = "Студент группы ИТ-101\nИванов Иван Иванович"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Слайд 2: Общая информация
    bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    
    title2 = slide2.shapes.title
    title2.text = "Общая информация"
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Имя: Петров Петр"
    
    p = tf.add_paragraph()
    p.text = "Возраст: 20 лет"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Место учебы: Воронежский государственный университет"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Специальность: Информационные технологии"
    p.level = 0
    
    # Слайд 3: Хобби и интересы
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    
    title3 = slide3.shapes.title
    title3.text = "Хобби и интересы"
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Программирование"
    
    interests = [
        "Python и веб-разработка",
        "Машинное обучение и AI",
        "Open Source проекты"
    ]
    
    for interest in interests:
        p = tf3.add_paragraph()
        p.text = interest
        p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Спорт"
    p.level = 0
    
    sports = [
        "Футбол",
        "Плавание",
        "Бег по утрам"
    ]
    
    for sport in sports:
        p = tf3.add_paragraph()
        p.text = sport
        p.level = 1
    
    # Слайд 4: Достижения
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    
    title4 = slide4.shapes.title
    title4.text = "Достижения"
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Академические"
    
    achievements = [
        "Диплом 1 степени на олимпиаде по программированию",
        "Средний балл 4.8",
        "Публикация статьи в научном журнале"
    ]
    
    for ach in achievements:
        p = tf4.add_paragraph()
        p.text = ach
        p.level = 1
    
    p = tf4.add_paragraph()
    p.text = "Профессиональные"
    p.level = 0
    
    prof_ach = [
        "Стажировка в IT-компании",
        "Участие в хакатонах",
        "Разработка собственных проектов"
    ]
    
    for ach in prof_ach:
        p = tf4.add_paragraph()
        p.text = ach
        p.level = 1
    
    # Слайд 5: Планы на будущее
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    
    title5 = slide5.shapes.title
    title5.text = "Планы на будущее"
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Краткосрочные цели (1 год)"
    
    short_goals = [
        "Завершить обучение с отличием",
        "Получить сертификаты по Python и Cloud",
        "Найти работу в крупной IT-компании"
    ]
    
    for goal in short_goals:
        p = tf5.add_paragraph()
        p.text = goal
        p.level = 1
    
    p = tf5.add_paragraph()
    p.text = "Долгосрочные цели (5 лет)"
    p.level = 0
    
    long_goals = [
        "Стать Team Lead или Architect",
        "Получить степень магистра",
        "Запустить собственный стартап"
    ]
    
    for goal in long_goals:
        p = tf5.add_paragraph()
        p.text = goal
        p.level = 1
    
    # Слайд 6: Контакты
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    
    title6 = slide6.shapes.title
    title6.text = "Контактная информация"
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Email: petrov.petr@example.com"
    
    contacts = [
        "Telegram: @petrov_petr",
        "GitHub: github.com/petrov-petr",
        "LinkedIn: linkedin.com/in/petrov-petr"
    ]
    
    for contact in contacts:
        p = tf6.add_paragraph()
        p.text = contact
        p.level = 0
    
    # Слайд 7: Заключение
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Добавляем текстовое поле для заключения
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    textbox = slide7.shapes.add_textbox(left, top, width, height)
    tf7 = textbox.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Сохранение презентации
    filename = 'labs/Lab8_PowerPoint_Presentation.pptx'
    prs.save(filename)
    print(f"✓ Лабораторная работа №8 успешно создана: {filename}")
    print(f"  Количество слайдов: {len(prs.slides)}")
    print(f"  Темы слайдов:")
    print(f"    1. Титульный слайд")
    print(f"    2. Общая информация")
    print(f"    3. Хобби и интересы")
    print(f"    4. Достижения")
    print(f"    5. Планы на будущее")
    print(f"    6. Контактная информация")
    print(f"    7. Заключение")
    return filename

if __name__ == "__main__":
    create_lab8_presentation()
